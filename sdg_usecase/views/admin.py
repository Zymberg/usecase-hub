import streamlit as st
import json
from utils.database import get_all_usecases, delete_usecase

def render():
    role = st.session_state.get("user_role", "Viewer")
    if role != "Admin":
        st.warning("⛔ Admin access required. Switch your role to Admin in the sidebar.")
        return

    st.markdown('<div class="page-title">Admin Panel</div>', unsafe_allow_html=True)
    st.markdown('<div class="page-subtitle">Manage all use cases and import from Google Drive.</div>', unsafe_allow_html=True)

    all_uc = get_all_usecases()

    # ── Stats ──────────────────────────────────────────────────────────────
    st.markdown("### Overview")
    c1, c2, c3 = st.columns(3)
    for col, num, label in [
        (c1, len(all_uc), "Total"),
        (c2, sum(1 for u in all_uc if u.get("status") == "Completed"), "Completed"),
        (c3, sum(1 for u in all_uc if u.get("status") == "In Progress"), "In Progress"),
    ]:
        col.markdown(f'<div class="stat-card"><div class="stat-num">{num}</div><div class="stat-label">{label}</div></div>', unsafe_allow_html=True)

    st.markdown("---")

    # ── Google Drive + OpenAI Sync ─────────────────────────────────────────
    st.markdown("### 📁 Import from Google Drive via OpenAI")
    st.info(
        "Each **subfolder** becomes one card (named after the folder, all files inside are read together). "
        "Each **loose file** in the root becomes its own card (named after the file). "
        "Every item gets a card — PDFs, images, PowerPoints, Docs, etc."
    )

    # Persist values in session state so they survive button clicks
    if "gdrive_folder_id"  not in st.session_state: st.session_state["gdrive_folder_id"]  = ""
    if "gdrive_creds_file" not in st.session_state: st.session_state["gdrive_creds_file"] = ""
    if "gdrive_openai_key" not in st.session_state: st.session_state["gdrive_openai_key"] = ""

    with st.expander("⚙️ Configure Google Drive + OpenAI"):
        folder_id  = st.text_input("Google Drive Folder ID",
                                   value=st.session_state["gdrive_folder_id"],
                                   placeholder="1NSbo21_itG1zO19ZsTg5jkX9q6L7XRki",
                                   help="From the folder URL: drive.google.com/drive/folders/<ID>",
                                   key="gdrive_folder_id")
        creds_file = st.text_input("Path to credentials.json",
                                   value=st.session_state["gdrive_creds_file"],
                                   placeholder="C:/Users/granit.zymeri/Desktop/sdg_usecase/utils/credentials.json",
                                   help="Google service account credentials JSON file",
                                   key="gdrive_creds_file")
        openai_key = st.text_input("OpenAI API Key",
                                   value=st.session_state["gdrive_openai_key"],
                                   type="password",
                                   placeholder="sk-proj-...",
                                   help="Uses gpt-4o — same key as your .env file",
                                   key="gdrive_openai_key")

        st.caption(
            "Subfolders → one card per folder (all files inside combined). "
            "Loose files → one card per file. "
            "Already-imported items won't be duplicated."
        )

        col_sync, col_debug = st.columns([1, 1])
        with col_sync:
            if st.button("🤖  Sync with OpenAI", use_container_width=True):
                _fid  = st.session_state.get("gdrive_folder_id", "").strip()
                _cred = st.session_state.get("gdrive_creds_file", "").strip()
                _key  = st.session_state.get("gdrive_openai_key", "").strip()
                if not _fid or not _cred or not _key:
                    st.error("Please fill in all three fields above.")
                else:
                    try:
                        _run_openai_sync(_fid, _cred, _key)
                    except Exception as e:
                        st.error(f"Sync failed: {e}")
        with col_debug:
            if st.button("🔍  Test Connection", use_container_width=True):
                _fid  = st.session_state.get("gdrive_folder_id", "").strip()
                _cred = st.session_state.get("gdrive_creds_file", "").strip()
                if not _fid or not _cred:
                    st.error("Please fill in Folder ID and credentials path.")
                else:
                    _test_drive_connection(_fid, _cred)

    st.markdown("""
**Setup steps:**
1. [Google Cloud Console](https://console.cloud.google.com/) → Enable **Google Drive API**
2. Create a Service Account → Download `credentials.json`
3. Share your Drive folder with the service account email
4. Install deps: `pip install google-api-python-client google-auth pdfminer.six python-pptx python-docx openpyxl`
    """)

    st.markdown("---")

    # ── Danger zone ────────────────────────────────────────────────────────
    st.markdown("### 🗑️ Danger Zone")
    col_del, col_confirm = st.columns([1, 3])
    with col_del:
        if st.button("🗑️ Delete All Use Cases", use_container_width=True):
            st.session_state["confirm_delete_all"] = True
    if st.session_state.get("confirm_delete_all"):
        st.warning("⚠️ This will permanently delete **all** use cases. This cannot be undone.")
        col_yes, col_no, _ = st.columns([1, 1, 3])
        with col_yes:
            if st.button("Yes, delete all", key="confirm_yes", use_container_width=True):
                from utils.database import get_conn
                conn = get_conn()
                conn.execute("DELETE FROM usecases")
                conn.commit()
                conn.close()
                st.session_state["confirm_delete_all"] = False
                st.success("✅ All use cases deleted.")
                st.rerun()
        with col_no:
            if st.button("Cancel", key="confirm_no", use_container_width=True):
                st.session_state["confirm_delete_all"] = False
                st.rerun()

    st.markdown("---")

    # ── All use cases table ────────────────────────────────────────────────
    st.markdown("### 📋 All Use Cases")
    search_q = st.text_input("", placeholder="Filter by title, team, client…", label_visibility="collapsed")
    filtered = all_uc
    if search_q:
        q = search_q.lower()
        filtered = [u for u in all_uc if q in (
            u.get("title","") + u.get("team_name","") + u.get("client","")
        ).lower()]

    if not filtered:
        st.info("No use cases found.")
        return

    for uc in filtered:
        with st.expander(f"**{uc['title']}** — {uc.get('team_name','?')} · {uc.get('status','?')} · {uc.get('created_at','')[:10]}"):
            c1, c2, c3 = st.columns([2, 2, 1])
            with c1:
                st.write(f"**Client:** {'Anonymized' if uc.get('client_anon') else uc.get('client','—')}")
                st.write(f"**Category:** {uc.get('category','—')}")
                st.write(f"**Source:** {uc.get('source','manual')}")
            with c2:
                st.write(f"**Team:** {uc.get('team_members','—')}")
                st.write(f"**Dates:** {uc.get('start_date','?')} → {uc.get('end_date','?')}")
            with c3:
                if st.button("✏️ Edit", key=f"admin_edit_{uc['id']}", use_container_width=True):
                    st.session_state["edit_id"] = uc["id"]
                    st.session_state["page"] = "submit"
                    st.rerun()
                if st.button("🗑️ Delete", key=f"admin_del_{uc['id']}", use_container_width=True):
                    delete_usecase(uc["id"])
                    st.success(f"Deleted: {uc['title']}")
                    st.rerun()


# ── Sync logic ─────────────────────────────────────────────────────────────

def _run_openai_sync(folder_id: str, creds_path: str, openai_key: str):
    from googleapiclient.discovery import build
    from google.oauth2 import service_account
    import io, json as _json, base64, requests
    from utils.database import save_usecase, get_conn

    # Auth
    creds = service_account.Credentials.from_service_account_file(
        creds_path, scopes=["https://www.googleapis.com/auth/drive.readonly"]
    )
    drive = build("drive", "v3", credentials=creds)

    st.info("🔍 Scanning Google Drive folder structure…")

    # Get root folder name
    root_meta = drive.files().get(fileId=folder_id, fields="name").execute()
    root_name = root_meta.get("name", "Root Folder")

    # Build a dict: { folder_name: { id, files: [...] } }
    folders = _list_folders_with_files(drive, folder_id, root_name)

    if not folders:
        st.warning("No subfolders or files found.")
        return

    st.info(f"Found **{len(folders)}** folder(s) to process: {', '.join(f'`{n}`' for n in folders)}")

    # Avoid duplicates — check ALL existing titles regardless of source
    conn = get_conn()
    existing = {r[0].strip().lower() for r in conn.execute("SELECT title FROM usecases").fetchall()}
    conn.close()

    progress   = st.progress(0)
    status_txt = st.empty()
    imported = skipped = errors = 0
    folder_list = list(folders.items())

    for idx, (folder_name, folder_data) in enumerate(folder_list):
        progress.progress((idx + 1) / len(folder_list))

        if folder_name.strip().lower() in existing:
            status_txt.caption(f"⏭️ Already exists: {folder_name}")
            skipped += 1
            continue

        files = folder_data["files"]
        folder_drive_id = folder_data["id"]
        status_txt.caption(f"📂 Processing folder: **{folder_name}** ({len(files)} file(s))")

        if not files:
            skipped += 1
            continue

        try:
            # Collect content from all files in this folder
            text_parts  = []
            image_parts = []

            for f in files:
                fname = f["name"]
                fid   = f["id"]
                fmime = f["mimeType"]
                from pathlib import Path
                ext = Path(fname).suffix.upper() or "?"
                status_txt.caption(f"📂 {folder_name} → reading [{ext}] {fname}")

                file_bytes, actual_mime = _download_file(drive, fid, fmime, fname)
                if file_bytes is None:
                    continue

                if actual_mime in ("image/png", "image/jpeg", "image/gif", "image/webp"):
                    b64 = base64.b64encode(file_bytes).decode()
                    image_parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:{actual_mime};base64,{b64}", "detail": "high"}
                    })
                else:
                    text = _extract_text(file_bytes, actual_mime, fname)
                    if text.strip():
                        text_parts.append(f"--- {fname} ---\n{text}")

            if not text_parts and not image_parts:
                skipped += 1
                continue

            # Build prompt — folder name IS the use case title
            prompt = f"""You are analyzing a use case called "{folder_name}" from a consulting/data team's portfolio.
The folder contains {len(files)} file(s). Read all provided content and extract a structured summary.
Use "{folder_name}" as the title exactly.
Return valid JSON only — no markdown fences, no extra text:
{{
  "title": "{folder_name}",
  "category": "<one of: Recommendation Engine, Scoring & Analytics, NLP / Text Analytics, Predictive Analytics, Fraud & Risk, Dashboard & Reporting, Automation & Pipelines, Computer Vision, Generative AI, Data Engineering, Other>",
  "client": "<client name if mentioned, else empty string>",
  "team_name": "<team name if mentioned, else empty string>",
  "team_members": "<comma-separated names if mentioned, else empty string>",
  "description": "<2-3 sentence description of the problem or use case>",
  "solution": "<2-3 sentence description of the solution or approach>",
  "tech_stack": "<comma-separated technologies mentioned>",
  "outcome": "<key results or metrics if mentioned, else empty string>",
  "tags": ["<tag1>", "<tag2>", "<tag3>"]
}}"""

            # Compose message: text first, then images
            content = [{"type": "text", "text": prompt}]
            if text_parts:
                content.append({"type": "text", "text": "\n\n".join(text_parts)[:8000]})
            content.extend(image_parts[:5])  # max 5 images per call

            resp = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {openai_key}", "Content-Type": "application/json"},
                json={"model": "gpt-4o", "max_tokens": 1200,
                      "messages": [{"role": "user", "content": content}]},
                timeout=90,
            )
            resp.raise_for_status()
            raw = resp.json()["choices"][0]["message"]["content"]
            raw = raw.strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
            data = _json.loads(raw)

            # Clean up the title using GPT
            clean_title = _clean_title(folder_name, openai_key)

            # Build Drive link — open?id= works universally for folders and all file types
            item_type = folder_data.get("type", "folder")
            if item_type == "folder":
                folder_link = f"https://drive.google.com/open?id={folder_drive_id}"
            else:
                file_id = folder_data["files"][0]["id"]
                folder_link = f"https://drive.google.com/open?id={file_id}"

            save_usecase({
                "id":           None,
                "title":        clean_title,  # GPT-cleaned title
                "category":     data.get("category") or "Other",
                "client":       data.get("client") or "",
                "client_anon":  0,
                "team_name":    data.get("team_name") or "",
                "team_members": data.get("team_members") or "",
                "start_date":   "",
                "end_date":     "",
                "description":  data.get("description") or "",
                "solution":     data.get("solution") or "",
                "tech_stack":   data.get("tech_stack") or "",
                "outcome":      data.get("outcome") or "",
                "links":        [folder_link],
                "tags":         data.get("tags") or ["gdrive-import"],
                "status":       "Completed",
                "source":       "gdrive",
            })
            existing.add(clean_title.strip().lower())
            existing.add(folder_name.strip().lower())  # also block original name
            imported += 1

        except Exception as e:
            errors += 1
            st.warning(f"⚠️ Failed on `{folder_name}`: {e}")
            continue

    progress.empty()
    status_txt.empty()
    st.success(f"✅ Done: **{imported} imported** · {skipped} skipped · {errors} error(s)")
    if imported:
        st.info("💡 Check Browse Library to review the new cards. Click Edit on any card to fill in missing details.")


def _list_folders_with_files(service, root_folder_id: str, root_name: str) -> dict:
    """
    Returns a dict of { card_name: { id, files: [...] } }.
    Rules:
    - Each subfolder → one card (named after folder, reads all files inside)
    - Each loose file in root → one card (named after the file, without extension)
    This means every item gets its own card regardless of format.
    """
    result = {}

    page_token = None
    while True:
        resp = service.files().list(
            q=f"'{root_folder_id}' in parents and trashed=false",
            fields="nextPageToken, files(id, name, mimeType)",
            pageSize=100,
            pageToken=page_token,
        ).execute()

        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                # Subfolder → card named after folder, reads all files inside
                files = _get_files_in_folder(service, f["id"])
                result[f["name"]] = {"id": f["id"], "type": "folder", "files": files}
            else:
                # Loose file → card named after file (without extension)
                card_name = f["name"].rsplit(".", 1)[0] if "." in f["name"] else f["name"]
                # Avoid name collision with a folder
                while card_name in result:
                    card_name = card_name + " (file)"
                result[card_name] = {"id": f["id"], "type": "file", "files": [f]}

        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    return result


def _get_files_in_folder(service, folder_id: str) -> list:
    """Get all files (non-folders) directly inside a folder."""
    files = []
    page_token = None
    while True:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false and mimeType != 'application/vnd.google-apps.folder'",
            fields="nextPageToken, files(id, name, mimeType)",
            pageSize=100,
            pageToken=page_token,
        ).execute()
        files.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return files


def _detect_mime(drive_mime: str, filename: str) -> str:
    EXT_MAP = {
        ".pdf":  "application/pdf",
        ".png":  "image/png",
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif":  "image/gif",
        ".webp": "image/webp",
        ".bmp":  "image/png",
        ".tiff": "image/jpeg",
        ".tif":  "image/jpeg",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt":  "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc":  "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls":  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv":  "text/csv",
        ".txt":  "text/plain",
        ".md":   "text/plain",
    }
    if drive_mime.startswith("application/vnd.google-apps."):
        return drive_mime
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return EXT_MAP.get(ext, drive_mime)


def _download_file(service, fid: str, drive_mime: str, filename: str = ""):
    from googleapiclient.http import MediaIoBaseDownload
    import io

    mime = _detect_mime(drive_mime, filename)

    EXPORT_MAP = {
        "application/vnd.google-apps.document":     ("text/plain", "text/plain"),
        "application/vnd.google-apps.presentation": ("text/plain", "text/plain"),
        "application/vnd.google-apps.spreadsheet":  ("text/csv",   "text/csv"),
    }
    SUPPORTED = {
        "application/pdf", "text/plain", "text/csv",
        "image/png", "image/jpeg", "image/gif", "image/webp",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }

    buf = io.BytesIO()
    try:
        if mime in EXPORT_MAP:
            export_mime, return_mime = EXPORT_MAP[mime]
            req = service.files().export_media(fileId=fid, mimeType=export_mime)
        elif mime in SUPPORTED:
            req = service.files().get_media(fileId=fid)
            return_mime = mime
        else:
            return None, None

        downloader = MediaIoBaseDownload(buf, req)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        return buf.getvalue(), return_mime
    except Exception:
        return None, None


def _extract_text(file_bytes: bytes, mime: str, name: str) -> str:
    import io

    if mime in ("text/plain", "text/csv"):
        return file_bytes.decode("utf-8", errors="ignore")[:6000]

    if mime == "application/pdf":
        try:
            from pdfminer.high_level import extract_text as pdf_extract
            return (pdf_extract(io.BytesIO(file_bytes)) or "")[:6000]
        except Exception:
            return file_bytes.decode("utf-8", errors="ignore")[:6000]

    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            from pptx import Presentation
            parts = []
            for slide in Presentation(io.BytesIO(file_bytes)).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)[:6000]
        except Exception:
            return ""

    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        try:
            import docx as _docx
            doc = _docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:6000]
        except Exception:
            return ""

    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)[:6000]
        except Exception:
            return ""

    return ""


def _test_drive_connection(folder_id: str, creds_path: str):
    """Diagnose Drive connection issues step by step."""
    import os

    # Step 1: credentials file exists?
    st.markdown("**🔍 Connection Diagnostics**")
    if not os.path.exists(creds_path):
        st.error(f"❌ credentials.json not found at: `{creds_path}`\n\nCheck the path — use forward slashes e.g. `C:/Users/name/Desktop/file.json`")
        return
    st.success(f"✅ credentials.json found")

    # Step 2: parse credentials
    try:
        import json as _json
        with open(creds_path) as f:
            creds_data = _json.load(f)
        svc_email = creds_data.get("client_email", "unknown")
        st.success(f"✅ Credentials valid — service account: `{svc_email}`")
        st.info(f"👆 Make sure this email has been given **Viewer** access to your Drive folder.")
    except Exception as e:
        st.error(f"❌ Could not parse credentials.json: {e}")
        return

    # Step 3: connect to Drive API
    try:
        from googleapiclient.discovery import build
        from google.oauth2 import service_account
        creds = service_account.Credentials.from_service_account_file(
            creds_path, scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )
        drive = build("drive", "v3", credentials=creds)
        st.success("✅ Connected to Google Drive API")
    except Exception as e:
        st.error(f"❌ Drive API connection failed: {e}\n\nMake sure `google-api-python-client` and `google-auth` are installed.")
        return

    # Step 4: can we read the folder?
    try:
        meta = drive.files().get(fileId=folder_id, fields="name, id").execute()
        st.success(f"✅ Folder found: **{meta['name']}** (ID: `{meta['id']}`)")
    except Exception as e:
        st.error(
            f"❌ Cannot access folder ID `{folder_id}`: {e}\n\n"
            f"Most likely cause: the folder has **not been shared** with `{svc_email}`.\n\n"
            f"Fix: Right-click the folder in Google Drive → Share → add `{svc_email}` as Viewer."
        )
        return

    # Step 5: list contents
    try:
        resp = drive.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="files(id, name, mimeType)",
            pageSize=20,
        ).execute()
        items = resp.get("files", [])
        if not items:
            st.warning("⚠️ Folder is accessible but **empty** — no files or subfolders found.")
        else:
            st.success(f"✅ Folder contains **{len(items)}** item(s):")
            for item in items:
                is_folder = item["mimeType"] == "application/vnd.google-apps.folder"
                icon = "📁" if is_folder else "📄"
                if is_folder:
                    link = f"https://drive.google.com/open?id={item['id']}"
                else:
                    link = f"https://drive.google.com/open?id={item['id']}"
                st.markdown(f"&nbsp;&nbsp;{icon} [{item['name']}]({link}) — `{item['mimeType']}`")
    except Exception as e:
        st.error(f"❌ Could not list folder contents: {e}")
        return

    st.success("🎉 Everything looks good! You can now run the full sync.")


def _clean_title(raw_name: str, openai_key: str) -> str:
    """
    Use GPT to turn a messy file/folder name into a clean use case title.
    Examples:
      "2024 SDG Use Case for Netun. Value proposition.ptx" → "2024 SDG Use Case - Netun"
      "Perfect Store"                                       → "Perfect Store"
      "Action Rec System v3 FINAL (2)"                     → "Action Recommendation System"
    Falls back to the raw name if GPT fails or key is missing.
    """
    import requests as _req

    if not openai_key:
        return raw_name

    try:
        resp = _req.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {openai_key}", "Content-Type": "application/json"},
            json={
                "model": "gpt-4o",
                "max_tokens": 30,
                "temperature": 0,
                "messages": [{
                    "role": "user",
                    "content": (
                        f"Convert this file or folder name into a clean, concise use case title. "
                        f"Rules: remove file extensions, version numbers (v1/v2/FINAL), "
                        f"generic words ('value proposition', 'deck', 'presentation', 'document'), "
                        f"extra punctuation and parentheses. "
                        f"If a client or company name is mentioned, keep it after a dash. "
                        f"Return ONLY the cleaned title, nothing else.\n\n"
                        f"Name: {raw_name}"
                    )
                }]
            },
            timeout=15,
        )
        resp.raise_for_status()
        title = resp.json()["choices"][0]["message"]["content"].strip().strip('"').strip("'")
        return title if title else raw_name
    except Exception:
        return raw_name
